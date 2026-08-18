#!/usr/bin/env python3
"""Lesson JSON -> .pptx deck.

Reads a lesson file from lessons/<name>.json and writes build/<name>.pptx.
The .pptx is designed to convert cleanly into native Google Slides on upload.

Usage:
    python generate_deck.py the-verb
    python generate_deck.py lessons/the-verb.json
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR

from themes import THEMES, get_theme

# ---- Theme (colors + brand set by apply_theme before a build) --------------
INK = ACCENT = MUTED = PAPER = RULE = ANSWER = DIAGRAM = None
EYEBROW = MOTTO = ""
HEAD_FONT = BODY_FONT = "Georgia"


def apply_theme(theme):
    global INK, ACCENT, MUTED, PAPER, RULE, ANSWER, DIAGRAM
    global EYEBROW, MOTTO, HEAD_FONT, BODY_FONT
    INK = RGBColor.from_string(theme["ink"])
    ACCENT = RGBColor.from_string(theme["accent"])
    MUTED = RGBColor.from_string(theme["muted"])
    PAPER = RGBColor.from_string(theme["paper"])
    RULE = RGBColor.from_string(theme["rule"])
    ANSWER = RGBColor.from_string(theme["answer"])
    DIAGRAM = RGBColor.from_string(theme["diagram"])
    EYEBROW = theme.get("eyebrow", "")
    MOTTO = theme.get("motto", "")
    HEAD_FONT = theme.get("head_font", "Georgia")
    BODY_FONT = theme.get("body_font", "Georgia")

EMU = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.9)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _run(p, text, *, size, color=None, bold=False, italic=False, font=None):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font if font is not None else BODY_FONT
    f.color.rgb = color if color is not None else INK
    return r


def _hanging_indent(p, amount):
    """Give a paragraph a hanging indent (wrapped lines indent by `amount` EMU)."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(amount)))
    pPr.set("indent", str(-int(amount)))


def _rule(slide, top, left=MARGIN, width=None, thickness=Pt(2), color=None):
    if width is None:
        width = SLIDE_W - 2 * MARGIN
    ln = slide.shapes.add_shape(1, left, top, width, thickness)  # rectangle
    ln.fill.solid()
    ln.fill.fore_color.rgb = color if color is not None else RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def _double_rule(slide, top, left=MARGIN, width=None):
    """Thick-then-thin letterpress rule, echoing Harvey's title pages."""
    _rule(slide, top, left=left, width=width, thickness=Pt(3), color=ACCENT)
    _rule(slide, top + Pt(5), left=left, width=width, thickness=Pt(1), color=RULE)


def _line(slide, x1, y1, x2, y2, color=None, width=Pt(1.5)):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color if color is not None else DIAGRAM
    cn.line.width = width
    cn.shadow.inherit = False
    return cn


# ---- Reed-Kellogg sentence diagram -----------------------------------------
DIAG_CHW = int(Inches(0.115))   # approx width per character at 18pt Georgia
DIAG_PAD = int(Inches(0.17))    # padding on each side of a word (gap for dividers)


def _slot_w(text):
    return len(text) * DIAG_CHW + 2 * DIAG_PAD


def _diag_word(slide, left, base_y, width, text):
    h = Inches(0.32)
    tb = slide.shapes.add_textbox(left, base_y - h, width, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, size=18, color=DIAGRAM)


def _draw_diagram(slide, left, top, s, p, o=None):
    """Baseline with subject | predicate | object.
    The subject/predicate divider crosses below the baseline; the
    predicate/object divider stops at it. Returns the right edge (EMU)."""
    base = top + Inches(0.30)
    top_line = top + Inches(0.04)
    below = Inches(0.13)
    ws, wp = _slot_w(s), _slot_w(p)
    wo = _slot_w(o) if o else 0
    x1 = left + ws
    x2 = x1 + wp
    xend = x2 + wo
    _line(slide, left, base, xend, base)               # baseline
    _line(slide, x1, top_line, x1, base + below)        # subject | predicate (crosses)
    if o:
        _line(slide, x2, top_line, x2, base)            # predicate | object (stops)
    _diag_word(slide, left, base, ws, s)
    _diag_word(slide, x1, base, wp, p)
    if o:
        _diag_word(slide, x2, base, wo, o)
    return xend


def _footer(slide, source):
    _, tf = _box(slide, MARGIN, SLIDE_H - Inches(0.6),
                 SLIDE_W - 2 * MARGIN, Inches(0.4))
    p = tf.paragraphs[0]
    _run(p, source, size=10, color=MUTED, italic=True)


def _heading(slide, text):
    _, tf = _box(slide, MARGIN, Inches(0.7),
                 SLIDE_W - 2 * MARGIN, Inches(1.0))
    p = tf.paragraphs[0]
    _run(p, text, size=34, color=ACCENT, bold=True, font=HEAD_FONT)
    _double_rule(slide, Inches(1.65))


# ---- Slide renderers -------------------------------------------------------
def render_title(slide, s, source):
    _bg(slide)
    # Eyebrow, small-caps, echoing the book's title page
    _, tfe = _box(slide, MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(0.5))
    pe = tfe.paragraphs[0]
    pe.alignment = PP_ALIGN.CENTER
    _run(pe, EYEBROW, size=15, color=MUTED, bold=True, font=HEAD_FONT)

    _double_rule(slide, Inches(2.45), left=Inches(4.4), width=Inches(4.5))

    _, tf = _box(slide, MARGIN, Inches(2.7), SLIDE_W - 2 * MARGIN, Inches(2.2))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, s["title"], size=60, color=ACCENT, bold=True, font=HEAD_FONT)
    if s.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(18)
        _run(p2, s["subtitle"], size=22, color=MUTED, italic=True)

    _double_rule(slide, Inches(5.05), left=Inches(4.4), width=Inches(4.5))

    if MOTTO:
        _, tfm = _box(slide, MARGIN, Inches(5.4), SLIDE_W - 2 * MARGIN, Inches(0.5))
        pm = tfm.paragraphs[0]
        pm.alignment = PP_ALIGN.CENTER
        _run(pm, MOTTO, size=15, color=ACCENT, italic=True)

    _footer(slide, source)


def render_definition(slide, s, source):
    _bg(slide)
    _heading(slide, s["term"])
    _, tf = _box(slide, MARGIN, Inches(2.1),
                 SLIDE_W - 2 * MARGIN, Inches(1.6))
    p = tf.paragraphs[0]
    _run(p, s["text"], size=30, color=INK, italic=True, font=BODY_FONT)

    if s.get("examples"):
        _, tf2 = _box(slide, MARGIN, Inches(3.9),
                      SLIDE_W - 2 * MARGIN, Inches(2.6))
        head = tf2.paragraphs[0]
        _run(head, "EXAMPLES", size=15, color=ACCENT, bold=True)
        for ex in s["examples"]:
            p = tf2.add_paragraph()
            p.space_before = Pt(6)
            _run(p, "•  ", size=22, color=ACCENT, bold=True)
            _run(p, ex, size=22, color=INK)
    _footer(slide, source)


def render_concept(slide, s, source):
    _bg(slide)
    _heading(slide, s["heading"])
    _, tf = _box(slide, MARGIN, Inches(2.0),
                 SLIDE_W - 2 * MARGIN, Inches(4.6))
    first = True
    for b in s["bullets"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        _run(p, "•  ", size=24, color=ACCENT, bold=True)
        _run(p, b, size=24, color=INK)
    _footer(slide, source)


def render_practice(slide, s, source, show_answers=False):
    _bg(slide)
    _heading(slide, s["heading"])
    if show_answers:
        _, tft = _box(slide, SLIDE_W - MARGIN - Inches(2.6), Inches(0.92),
                      Inches(2.6), Inches(0.4))
        pt = tft.paragraphs[0]
        pt.alignment = PP_ALIGN.RIGHT
        _run(pt, "ANSWER KEY", size=14, color=ACCENT, bold=True)

    _, tf = _box(slide, MARGIN, Inches(1.95),
                 SLIDE_W - 2 * MARGIN, Inches(1.1))
    p = tf.paragraphs[0]
    _run(p, s["instruction"], size=18, color=MUTED, italic=True)

    answers = s.get("answers") or []
    is_diagram = bool(answers) and isinstance(answers[0], dict)

    if show_answers and is_diagram:
        _practice_diagram_answers(slide, s, answers)
    else:
        _practice_item_columns(slide, s, answers if show_answers else None)
    _footer(slide, source)


def _practice_item_columns(slide, s, answers):
    """Two columns of numbered items. Answers sit on their own line beneath
    each item, unless the slide sets "answerInline" (then they follow inline —
    good for long sentences with a one-word answer)."""
    items = s["items"]
    inline = s.get("answerInline", False)
    mid = (len(items) + 1) // 2
    cols = [items[:mid], items[mid:]]
    gap = Inches(0.3)                         # tighter gap between columns
    usable = SLIDE_W - MARGIN - Inches(0.6)   # narrower right cushion -> wider columns
    col_w = (usable - gap) / 2
    for ci, col in enumerate(cols):
        left = MARGIN + ci * (col_w + gap)
        _, tf = _box(slide, left, Inches(3.1), col_w, Inches(3.5))
        first = True
        for i, it in enumerate(col):
            n = ci * mid + i + 1
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10 if (not answers or inline) else 4)
            _run(p, f"{n}.  ", size=22, color=ACCENT, bold=True)
            _run(p, it, size=22, color=INK)
            if answers and n - 1 < len(answers):
                if inline:
                    _run(p, "   ", size=22, color=ANSWER)
                    _run(p, answers[n - 1], size=22, color=ANSWER, bold=True)
                else:
                    pa = tf.add_paragraph()
                    pa.space_after = Pt(14)
                    _run(pa, "     ", size=20, color=ANSWER)   # small hanging indent
                    _run(pa, answers[n - 1], size=20, color=ANSWER, bold=True)


def _practice_diagram_answers(slide, s, answers):
    """Single column of Reed-Kellogg diagrams, one per item."""
    items = s["items"]
    y = Inches(3.15)
    row_h = Inches(0.72)
    for i, a in enumerate(answers):
        _, tf = _box(slide, MARGIN, y - Inches(0.02), Inches(0.5), Inches(0.4))
        _run(tf.paragraphs[0], f"{i + 1}.", size=18, color=ACCENT, bold=True)
        _draw_diagram(slide, MARGIN + Inches(0.55), y,
                      a["s"], a["p"], a.get("o"))
        y += row_h


def render_questions(slide, s, source, show_answers=False):
    _bg(slide)
    _heading(slide, "Questions")
    if show_answers:
        _, tft = _box(slide, SLIDE_W - MARGIN - Inches(2.6), Inches(0.92),
                      Inches(2.6), Inches(0.4))
        pt = tft.paragraphs[0]
        pt.alignment = PP_ALIGN.RIGHT
        _run(pt, "ANSWER KEY", size=14, color=ACCENT, bold=True)

    answers = s.get("answers") or []
    # The answer version packs a lot in, so it starts higher and sits tighter.
    top = Inches(1.75) if show_answers else Inches(1.9)
    _, tf = _box(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(5.2))
    first = True
    for i, q in enumerate(s["items"], 1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4 if show_answers else 14)
        if show_answers:
            p.line_spacing = 0.9
        _run(p, f"{i}.  ", size=22, color=ACCENT, bold=True)
        _run(p, q, size=22, color=INK)
        # Answer runs inline on the same line; hanging indent keeps a wrap aligned.
        if show_answers and i - 1 < len(answers):
            _hanging_indent(p, int(Inches(0.42)))
            _run(p, "   ", size=20, color=INK)
            _run(p, answers[i - 1], size=20, color=ANSWER, bold=True)
    _footer(slide, source)


RENDERERS = {
    "title": render_title,
    "definition": render_definition,
    "concept": render_concept,
    "practice": render_practice,
    "questions": render_questions,
}


def build(lesson_path: Path, theme_name="light") -> Path:
    apply_theme(get_theme(theme_name))
    data = json.loads(lesson_path.read_text())
    source = data.get("source", "")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in data["slides"]:
        slide = _blank(prs)
        RENDERERS[s["type"]](slide, s, source)
        # A practice/questions slide with answers is followed by an "Answer Key" slide.
        if s["type"] == "practice" and s.get("answers"):
            render_practice(_blank(prs), s, source, show_answers=True)
        elif s["type"] == "questions" and s.get("answers"):
            render_questions(_blank(prs), s, source, show_answers=True)

    out_dir = Path(__file__).resolve().parent / "build"
    out_dir.mkdir(exist_ok=True)
    suffix = "" if theme_name == "light" else f"-{theme_name}"
    out_path = out_dir / f"{data['lesson']}{suffix}.pptx"
    prs.save(str(out_path))
    return out_path


def resolve(arg: str) -> Path:
    p = Path(arg)
    if p.suffix == ".json" and p.exists():
        return p
    cand = Path(__file__).resolve().parents[3] / "content" / "grammar" / "lessons" / f"{arg}.json"
    if cand.exists():
        return cand
    raise SystemExit(f"Lesson not found: {arg}")


if __name__ == "__main__":
    lesson, theme = "the-verb", "light"
    for a in sys.argv[1:]:
        if a in THEMES:
            theme = a
        else:
            lesson = a
    out = build(resolve(lesson), theme)
    print(f"wrote {out}  (theme: {theme})")
